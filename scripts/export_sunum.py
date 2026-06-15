from __future__ import annotations

import argparse
import subprocess
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
LINE = RGBColor(209, 213, 219)
SOFT = RGBColor(246, 247, 249)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(180, 83, 9)
WHITE = RGBColor(255, 255, 255)


PDF_EXPORT_CSS = """

    @page {
      size: 13.333333in 7.5in;
      margin: 0;
    }

    @media print {
      html,
      body {
        width: 13.333333in;
        margin: 0;
        background: #ffffff;
        overflow: visible;
        -webkit-print-color-adjust: exact;
        print-color-adjust: exact;
      }

      body.pdf-export .deck {
        display: block;
        width: 13.333333in;
        height: auto;
        padding: 0;
      }

      body.pdf-export .slide {
        display: flex !important;
        width: 13.333333in;
        height: 7.5in;
        min-height: 0;
        aspect-ratio: auto;
        border: 0;
        box-shadow: none;
        padding: 52px 62px 44px;
        flex-direction: column;
        gap: 30px;
        break-after: page;
        page-break-after: always;
      }

      body.pdf-export .slide.compact-slide {
        gap: 14px;
        padding-top: 34px;
        padding-bottom: 26px;
      }

      body.pdf-export .controls {
        display: none !important;
      }

      body.pdf-export .two-col {
        grid-template-columns: 1fr 1fr !important;
      }

      body.pdf-export .split-visual {
        grid-template-columns: 0.98fr 1.02fr !important;
      }

      body.pdf-export .figure-layout {
        grid-template-columns: 1.14fr 0.86fr !important;
      }

      body.pdf-export .slide[data-index="3"] .metric-row {
        grid-template-columns: repeat(3, 1fr) !important;
      }

      body.pdf-export .slide[data-index="5"] .metric-row {
        grid-template-columns: 1fr !important;
      }

      body.pdf-export .process {
        grid-template-columns: repeat(7, 1fr) !important;
      }

      body.pdf-export .branch-row,
      body.pdf-export .flow-branches {
        grid-template-columns: 1fr 1fr !important;
      }

      body.pdf-export .flow-top {
        grid-template-columns: repeat(4, 1fr) !important;
      }

      body.pdf-export .flow-bottom {
        grid-template-columns: none !important;
      }

      body.pdf-export h1,
      body.pdf-export .title-slide h1 {
        font-size: 58px !important;
      }

      body.pdf-export h2 {
        font-size: 48px !important;
      }

      body.pdf-export .kicker {
        font-size: 24px !important;
      }

      body.pdf-export p,
      body.pdf-export ul,
      body.pdf-export ol {
        font-size: 30px !important;
      }

      body.pdf-export .metric .label {
        font-size: 22px !important;
      }

      body.pdf-export .flow-node .text,
      body.pdf-export .flow-branch p {
        font-size: 15px !important;
      }
    }
"""


def prepare_pdf_html(source: Path, destination: Path) -> None:
    content = source.read_text(encoding="utf-8")
    base_href = source.parent.resolve().as_uri() + "/"
    content = content.replace("<head>", f'<head>\n  <base href="{base_href}">', 1)
    content = content.replace("</style>", f"{PDF_EXPORT_CSS}\n  </style>", 1)
    content = content.replace("<body>", '<body class="pdf-export">', 1)
    destination.write_text(content, encoding="utf-8")


def make_pdf(chrome: str, html_path: Path, output_pdf: Path) -> None:
    with tempfile.TemporaryDirectory() as tmp:
        export_html = Path(tmp) / "sunum_pdf_export.html"
        prepare_pdf_html(html_path, export_html)
        command = [
            chrome,
            "--headless=new",
            "--no-sandbox",
            "--disable-gpu",
            "--no-pdf-header-footer",
            f"--print-to-pdf={output_pdf}",
            export_html.as_uri(),
        ]
        subprocess.run(command, check=True)


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    *,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def add_multiline_text(
    slide,
    lines: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    *,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        if align is not None:
            paragraph.alignment = align


def add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 25,
    color: RGBColor = INK,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(8)
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color


def add_footer(slide, index: int, total: int) -> None:
    add_text(slide, f"{index} / {total}", 12.45, 7.18, 0.55, 0.2, 12, color=MUTED, align=PP_ALIGN.RIGHT)


def add_kicker(slide, text: str, index: int, total: int = 8) -> None:
    add_text(slide, text, 0.65, 0.42, 8.0, 0.32, 19, color=BLUE, bold=True)
    add_footer(slide, index, total)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    accent: RGBColor = BLUE,
    accent_side: str = "left",
    fill: RGBColor = SOFT,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    if accent_side == "top":
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.07))
    else:
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()


def add_metric(slide, left: float, top: float, width: float, value: str, label: str, accent: RGBColor) -> None:
    add_card(slide, left, top, width, 1.18, accent=accent)
    add_text(slide, value, left + 0.28, top + 0.2, width - 0.42, 0.42, 27, bold=True)
    add_text(slide, label, left + 0.28, top + 0.7, width - 0.42, 0.3, 15, color=MUTED)


def add_callout(slide, text: str, left: float, top: float, width: float, height: float, *, accent: RGBColor = BLUE) -> None:
    add_card(slide, left, top, width, height, accent=accent)
    add_text(slide, text, left + 0.28, top + 0.28, width - 0.48, height - 0.45, 22, bold=True)


def add_picture(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    border.fill.solid()
    border.fill.fore_color.rgb = WHITE
    border.line.color.rgb = LINE
    border.line.width = Pt(1)

    with Image.open(image_path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        pic_width = width - 0.26
        pic_height = pic_width / ratio
    else:
        pic_height = height - 0.26
        pic_width = pic_height * ratio
    pic_left = left + (width - pic_width) / 2
    pic_top = top + (height - pic_height) / 2
    slide.shapes.add_picture(str(image_path), Inches(pic_left), Inches(pic_top), width=Inches(pic_width), height=Inches(pic_height))


def add_table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(15 if row_index else 14)
            paragraph.font.bold = row_index == 0
            paragraph.font.color.rgb = MUTED if row_index == 0 else INK
            paragraph.alignment = PP_ALIGN.CENTER if col_index else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE


def add_flow_node(slide, number: str, text: str, left: float, top: float, width: float, height: float, *, accent: RGBColor = BLUE) -> None:
    add_card(slide, left, top, width, height, accent=accent, accent_side="top")
    add_text(slide, number, left + 0.12, top + 0.15, 0.42, 0.25, 17, color=accent, bold=True)
    add_text(slide, text, left + 0.12, top + 0.47, width - 0.22, height - 0.52, 11)


def make_pptx(output_pptx: Path, figures_dir: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    blank = presentation.slide_layouts[6]
    presentation.core_properties.title = "LLM'lerde Matematiksel Problem Çözme"
    presentation.core_properties.author = "Muhammet Emin AYDINALP"

    # 1
    slide = presentation.slides.add_slide(blank)
    add_text(slide, "Hesaplamalı Anlambilim Final Çalışması", 3.9, 2.03, 5.6, 0.35, 18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "LLM'lerde Matematiksel Problem Çözme için Farklı LoRA İnce Ayar Yaklaşımlarının Değerlendirilmesi",
        1.25,
        2.65,
        10.85,
        1.25,
        34,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_multiline_text(slide, ["Muhammet Emin AYDINALP", "25501819"], 4.25, 4.75, 4.85, 0.7, 21, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 1, 8)

    # 2
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Problem ve Motivasyon", 2)
    add_bullets(
        slide,
        [
            "Matematiksel sözel problemler LLM'lerde akıl yürütmeyi ölçmek için kullanılır.",
            "Modelin nicelikleri ayırması, ilişkileri kurması ve hesap sırasını takip etmesi gerekir.",
            "Bu çalışmada öğretmen çözümüyle yapılan küçük güncellemelerin etkisi incelendi.",
        ],
        0.75,
        1.55,
        5.75,
        3.6,
        size=23,
    )
    add_callout(
        slide,
        "M1'in yanlış yaptığı bir soru öğretildiğinde, bu bilgi benzer sorulara ve bağımsız test kümesine aktarılabilir mi?",
        7.15,
        2.2,
        5.35,
        1.7,
    )

    # 3
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Veri ve Modeller", 3)
    add_metric(slide, 0.65, 1.25, 3.85, "GSM8K_TR", "Türkçe matematiksel sözel problem veri kümesi", BLUE)
    add_metric(slide, 4.78, 1.25, 3.85, "Qwen3.5-4B", "Temel model, çalışmada M1 olarak kullanıldı", GREEN)
    add_metric(slide, 8.91, 1.25, 3.85, "gpt-oss-120b", "Öğretmen çözümü ve benzer soru üretimi", AMBER)
    add_bullets(
        slide,
        [
            "Sabit test kümesi: 500 soru",
            "Eğitim kümesi: M1'in yanlış yaptığı ve öğretmen modelle doğrulanan 500 soru",
            "Eğitim ve test kümeleri arasında çakışma yoktur.",
        ],
        0.9,
        3.25,
        11.5,
        2.2,
        size=23,
    )

    # 4
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Deney Akışı", 4)
    top_y = 0.98
    box_w = 2.8
    gap = 0.28
    x_positions = [0.65 + i * (box_w + gap) for i in range(4)]
    nodes = [
        ("1", "M1 GSM8K_TR üzerinde değerlendirilir."),
        ("2", "Yanlış çözülen sorular belirlenir."),
        ("3", "Öğretmen modelden doğru çözüm alınır."),
        ("4", "Aynı yapıda benzer soru üretilir."),
    ]
    for x, (number, text) in zip(x_positions, nodes):
        add_flow_node(slide, number, text, x, top_y, box_w, 0.78)
    for x in [3.35, 6.43, 9.51]:
        add_text(slide, "→", x, top_y + 0.22, 0.25, 0.25, 18, color=MUTED, bold=True)
    add_text(slide, "↓", 6.43, 1.84, 0.25, 0.25, 18, color=MUTED, bold=True)
    add_flow_node(
        slide,
        "5",
        "Aktif adapter'dan aday adapter oluşturulur, tek soru-cevap çiftiyle 3 epoch eğitilir.",
        4.25,
        2.08,
        4.85,
        0.82,
        accent=AMBER,
    )
    add_text(slide, "↙", 5.15, 2.95, 0.28, 0.25, 18, color=MUTED, bold=True)
    add_text(slide, "↘", 7.88, 2.95, 0.28, 0.25, 18, color=MUTED, bold=True)
    add_card(slide, 2.25, 3.25, 4.35, 0.82, accent=BLUE)
    add_text(slide, "Selective kararı", 2.5, 3.4, 2.4, 0.24, 16, bold=True)
    add_text(slide, "Benzer soru doğruysa aday adapter kabul edilir. Yanlışsa güncelleme geri alınır.", 2.5, 3.72, 3.85, 0.25, 11)
    add_card(slide, 6.8, 3.25, 4.35, 0.82, accent=GREEN)
    add_text(slide, "Blind kararı", 7.05, 3.4, 2.4, 0.24, 16, bold=True)
    add_text(slide, "Benzer soru sonucu kaydedilir. Doğru ya da yanlış olsa da güncelleme tutulur.", 7.05, 3.72, 3.85, 0.25, 11)
    add_text(slide, "↘", 5.25, 4.15, 0.28, 0.25, 18, color=MUTED, bold=True)
    add_text(slide, "↙", 7.9, 4.15, 0.28, 0.25, 18, color=MUTED, bold=True)
    add_flow_node(slide, "6", "Benzer sorular üzerinde genel değerlendirme yapılır.", 4.22, 4.42, 4.9, 0.62)
    add_text(slide, "↓", 6.43, 5.16, 0.25, 0.25, 18, color=MUTED, bold=True)
    add_flow_node(slide, "7", "Bağımsız final test kümesinde karşılaştırma yapılır.", 4.22, 5.46, 4.9, 0.62)

    # 5
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Benzer Soru Başarısı", 5)
    add_picture(slide, figures_dir / "similar_learning_curve.png", 0.65, 1.45, 6.65, 4.25)
    add_metric(slide, 7.6, 1.1, 5.05, "405 / 500", "Selective doğru sayısı, doğruluk 0.810", BLUE)
    add_metric(slide, 7.6, 2.55, 5.05, "389 / 500", "Blind doğru sayısı, doğruluk 0.778", GREEN)
    add_callout(slide, "Benzer soru ölçütünde Selective, zararlı görünen güncellemeleri elediği için öne çıktı.", 7.6, 4.0, 5.05, 1.75)

    # 6
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Bağımsız Test Başarısı", 6)
    add_picture(slide, figures_dir / "test_accuracy.png", 0.65, 1.35, 6.5, 4.55)
    add_table(
        slide,
        [
            ["Model", "Doğruluk Oranı", "Doğruluk"],
            ["M1 baseline", "290 / 500", "0.580"],
            ["Selective", "274 / 500", "0.548"],
            ["Blind", "296 / 500", "0.592"],
        ],
        7.55,
        1.65,
        4.95,
        2.05,
    )

    # 7
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Düzeltme ve Bozma Dengesi", 7)
    add_picture(slide, figures_dir / "test_transitions.png", 0.65, 1.35, 6.55, 4.55)
    add_bullets(
        slide,
        [
            "Selective: 79 eski hatayı düzeltti, 95 eski doğruyu bozdu.",
            "Blind: 85 eski hatayı düzeltti, 79 eski doğruyu bozdu.",
            "Net etki: Selective -16, Blind +6.",
        ],
        7.55,
        1.65,
        4.85,
        2.65,
        size=23,
    )

    # 8
    slide = presentation.slides.add_slide(blank)
    add_kicker(slide, "Sonuç", 8)
    add_bullets(
        slide,
        [
            "Öğretmen çözümüyle yapılan tek örneklik LoRA güncellemesi benzer sorulara aktarım sağlayabiliyor.",
            "Benzer soru başarısı bağımsız test başarısını doğrudan garanti etmiyor.",
            "Selective yerel başarıda daha iyi, final testte daha zayıf kaldı.",
            "Blind daha basit olmasına rağmen bu deney koşullarında daha iyi genelledi.",
        ],
        0.75,
        1.2,
        11.9,
        2.55,
        size=22,
    )
    add_card(slide, 0.85, 4.25, 11.65, 1.82, accent=GREEN)
    add_text(slide, "Gelecek Çalışmalar", 1.15, 4.45, 5.0, 0.32, 21, bold=True)
    add_bullets(
        slide,
        [
            "Gelecek çalışmalarda genel performansı koruyacak güncelleme kabul kuralları geliştirilebilir.",
            "Daha çok veri seti ile ve farklı çözümler ile eğitim yapılabilir.",
            "Sadece sayısal cevap değil de çözümün tamamına odaklanan değerlendirme metrikleri kullanılabilir.",
        ],
        1.15,
        4.85,
        10.9,
        0.9,
        size=16,
    )

    presentation.save(output_pptx)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--html", default="sunum.html")
    parser.add_argument("--output-dir", default="results")
    parser.add_argument("--chrome", default="google-chrome")
    args = parser.parse_args()

    root = Path.cwd()
    html_path = (root / args.html).resolve()
    output_dir = (root / args.output_dir).resolve()
    figures_dir = output_dir / "figures"
    output_dir.mkdir(parents=True, exist_ok=True)

    pdf_path = output_dir / "sunum.pdf"
    pptx_path = output_dir / "sunum.pptx"

    make_pdf(args.chrome, html_path, pdf_path)
    make_pptx(pptx_path, figures_dir)

    print(pdf_path)
    print(pptx_path)


if __name__ == "__main__":
    main()
